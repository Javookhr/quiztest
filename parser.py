# ============================================================
#  parser.py — Fayllardan savollarni ajratib olish
#
#  Qo'llab-quvvatlanadigan formatlar:
#
#  FORMAT 1 — +++ bilan boshlanadigan savol (YANGI):
#      +++
#      Savol matni
#      ====
#      # To'g'ri javob    ← # belgisi to'g'ri javobni bildiradi
#      ====
#      Variant 2
#      ====
#      Variant 3
#      ====
#      Variant 4
#
#  FORMAT 2 — ==== ajratuvchi (klassik):
#      Savol matni
#      ====
#      # To'g'ri javob
#      ====
#      Variant 2
#      ====
#      Variant 3
#      ====
#      Variant 4
#
#  FORMAT 3 — Raqamli (klassik):
#      1. Savol matni?
#      A) Variant 1
#      B) Variant 2
#      C) To'g'ri variant
#      D) Variant 4
#      Javob: C
#
#  Fayllar: .txt  .docx  .pptx  .xlsx  .pdf  rasm(OCR)
# ============================================================

import io
import re
import random
import logging
from typing import List, Dict, Tuple

log = logging.getLogger(__name__)


# ──────────────────────────────────────────────────────────
#  Yordamchi: variantlarni aralashtirish
# ──────────────────────────────────────────────────────────
def shuffle_options(question: Dict) -> Dict:
    """
    Variantlarni aralashtiradi va correct_option_id ni yangilaydi.
    To'g'ri javob har safar turli o'rinda chiqadi.
    """
    opts         = question["options"][:]
    correct_text = opts[question["correct_option_id"]]

    random.shuffle(opts)
    new_correct = opts.index(correct_text)

    return {**question, "options": opts, "correct_option_id": new_correct}


# ──────────────────────────────────────────────────────────
#  FORMAT 1 — +++ va ==== bilan (yangi format)
# ──────────────────────────────────────────────────────────
def _parse_plus_equal_format(text: str) -> List[Dict]:
    """
    +++ yoki ++++ bilan boshlanuvchi savollarni o'qiydi.

    FORMAT A — ==== ajratuvchili (eski format):
        +++
        Savol matni
        ====
        # To'g'ri javob
        ====
        Variant 2
        ====
        Variant 3
        ====
        Variant 4

    FORMAT B — har bir qator alohida variant (yangi format):
        ++++
        Savol matni
        Variant 1
        Variant 2
        # To'g'ri javob
        Variant 4
    """
    # ++++ yoki +++ bilan ajratilgan bloklar
    blocks = re.split(r"\+{3,}", text)
    questions: List[Dict] = []

    for block in blocks:
        if not block.strip():
            continue

        # ==== ajratuvchi bormi? — FORMAT A
        if re.search(r"={3,}", block):
            # ====ni alohida qatorga normalize qilamiz ("savol?====" holatlarini tuzatish)
            normalized = re.sub(r"={3,}", "\n====\n", block)
            parts = normalized.split("\n====\n")
            # parts[0] = savol matni (====dan OLDIN kelgan qism)
            # Faqat optionlar uchun bo'sh qismlarni filtr qilamiz
            # savol qismi uchun ALOHIDA tekshiramiz

            if len(parts) < 2:
                continue

            # parts[0] bo'sh bo'lsa — savol yo'q (blok ====dan boshlangan), o'tkazib yub
            q_text = re.sub(r"^\++\s*", "", parts[0]).strip()
            q_text = re.sub(r"^[#\s]+", "", q_text).strip()
            if not q_text:
                continue

            option_parts = [p.strip() for p in parts[1:] if p.strip()]

            options: List[str] = []
            correct_id: int = 0

            for part in option_parts:
                if len(options) >= 4:
                    break
                part = part.strip()
                if not part:
                    continue
                if part.startswith("#"):
                    clean = part.lstrip("#").strip()
                    if clean:
                        correct_id = len(options)
                        options.append(clean)
                else:
                    options.append(part)

            if len(options) >= 2:
                questions.append({
                    "question": q_text,
                    "options": options[:4],
                    "correct_option_id": correct_id,
                })

        else:
            # FORMAT B — har bir qator alohida variant
            lines = [l.strip() for l in block.strip().splitlines() if l.strip()]
            # Blokning boshidagi ortiqcha + belgilerini olib tashlash
            lines = [re.sub(r"^\++\s*", "", l) for l in lines]
            lines = [l for l in lines if l]

            if len(lines) < 2:
                continue

            q_text = lines[0]
            q_text = re.sub(r"^[#\s]+", "", q_text).strip()
            options: List[str] = []
            correct_id: int = 0

            for line in lines[1:]:
                if len(options) >= 4:
                    break
                if line.startswith("#"):
                    clean = line.lstrip("#").strip()
                    correct_id = len(options)
                    options.append(clean)
                else:
                    options.append(line)

            if len(options) >= 2:
                questions.append({
                    "question": q_text,
                    "options": options[:4],
                    "correct_option_id": correct_id,
                })

    return [shuffle_options(q) for q in questions]


# ──────────────────────────────────────────────────────────
#  FORMAT 2 — ==== ajratuvchi bilan (klassik)
# ──────────────────────────────────────────────────────────
def _parse_equal_format(text: str) -> List[Dict]:
    """
    Quyidagi ko'rinishdagi faylni o'qiydi:

        Savol matni
        ====
        # To'g'ri javob
        ====
        Variant 2
        ====
        Variant 3
        ====
        Variant 4

    '# ' bilan boshlanuvchi variant — to'g'ri javob.
    Har bir savol + variantlar guruhi '====' bilan ajratiladi.
    """
    # ===, ====, ===== hammasini ajratuvchi sifatida qabul qilish
    parts = re.split(r"\n?={3,}\n?", text.strip())
    parts = [p.strip() for p in parts]

    questions: List[Dict] = []
    i = 0

    while i < len(parts):
        q_raw = parts[i]
        if not q_raw:
            i += 1
            continue

        # Raqamli prefixni olib tashlash (1. yoki 1) ko'rinishi)
        q_text = re.sub(r"^\d+[.)]\s*", "", q_raw).strip()
        q_text = re.sub(r"^[#\s]+", "", q_text).strip()
        if not q_text:
            i += 1
            continue

        # Keyingi qismlardan variantlarni to'plash
        options: List[str]  = []
        correct_id: int     = 0
        j = i + 1

        while j < len(parts) and len(options) < 4:
            raw_opt = parts[j].strip()
            if raw_opt:
                if raw_opt.startswith("#"):
                    clean = raw_opt.lstrip("#").strip()
                    correct_id = len(options)
                    options.append(clean)
                else:
                    options.append(raw_opt)
            j += 1

        if len(options) >= 2:
            questions.append({
                "question":        q_text,
                "options":         options[:4],
                "correct_option_id": correct_id,
            })

        i = j  # Keyingi savolga o'tish

    return [shuffle_options(q) for q in questions]


# ──────────────────────────────────────────────────────────
#  FORMAT 2 — Raqamli klassik format
# ──────────────────────────────────────────────────────────
_ANSWER_RE   = re.compile(
    r"(?:[Jj]avob|[Tt]o['\u2019]?g['\u2019]?ri(?:\s+javob)?|"
    r"[Аа]nswer|[Оо]тв(?:ет)?)\s*[:\-]?\s*([A-Da-d1-4])",
    re.UNICODE,
)
_QUESTION_RE = re.compile(r"^(\d+)[.)]\s*(.+)", re.UNICODE)
_OPTION_RE   = re.compile(r"^([A-Da-d1-4])[.)]\s*(.+)", re.UNICODE)


def _parse_classic_format(text: str) -> List[Dict]:
    """
    Klassik raqamli format:
        1. Savol?
        A) Variant
        B) Variant
        C) To'g'ri
        D) Variant
        Javob: C
    """
    questions: List[Dict] = []
    current_q: str | None = None
    current_opts: List[str] = []
    correct_id: int | None = None

    for raw_line in text.splitlines():
        line = raw_line.strip()
        if not line:
            continue

        # Savol qatori
        qm = _QUESTION_RE.match(line)
        if qm:
            if current_q and len(current_opts) >= 2 and correct_id is not None:
                questions.append({
                    "question":          current_q,
                    "options":           current_opts[:4],
                    "correct_option_id": correct_id,
                })
            current_q    = qm.group(2).strip()
            current_opts = []
            correct_id   = None
            continue

        # Variant qatori
        om = _OPTION_RE.match(line)
        if om:
            current_opts.append(om.group(2).strip())
            continue

        # Javob qatori
        am = _ANSWER_RE.match(line)
        if am:
            ans = am.group(1).upper()
            mapping = {"A": 0, "B": 1, "C": 2, "D": 3,
                       "1": 0, "2": 1, "3": 2, "4": 3}
            correct_id = mapping.get(ans, 0)

    # Oxirgi savolni saqlash
    if current_q and len(current_opts) >= 2 and correct_id is not None:
        questions.append({
            "question":          current_q,
            "options":           current_opts[:4],
            "correct_option_id": correct_id,
        })

    return [shuffle_options(q) for q in questions]


# ──────────────────────────────────────────────────────────
#  Asosiy matn parser — avval ==== formatni sinab ko'radi
# ──────────────────────────────────────────────────────────
def parse_questions_from_text(text: str) -> List[Dict]:
    """
    Uchta formatni ham qo'llab-quvvatlaydi:
    1) +++ va ==== ajratuvchi (yangi format)
    2) ==== ajratuvchi format
    3) Klassik 1. A) B) C) D) Javob: format
    """
    # +++ format bormi? — avval shuni sinab ko'r
    if "+++" in text:
        qs = _parse_plus_equal_format(text)
        if qs:
            return qs
    
    # ==== format bormi? — keyingi
    if re.search(r"={3,}", text):
        qs = _parse_equal_format(text)
        if qs:
            return qs

    # Klassik A) B) C) D) format
    return _parse_classic_format(text)


# ──────────────────────────────────────────────────────────
#  Fayl formatlari bo'yicha parserlar
# ──────────────────────────────────────────────────────────
def _parse_docx(data: bytes) -> str:
    from docx import Document
    try:
        doc   = Document(io.BytesIO(data))
        parts = [p.text for p in doc.paragraphs if p.text.strip()]
        for table in doc.tables:
            for row in table.rows:
                parts.append("\t".join(c.text.strip() for c in row.cells))
        return "\n".join(parts)
    except Exception as e:
        log.error(f"Docx o'qishda xatolik: {e}")
        raise ValueError(f"Word faylini o'qib bo'lmadi (docx formatida emas yoki buzilgan).")


def _parse_pptx(data: bytes) -> str:
    from pptx import Presentation
    prs   = Presentation(io.BytesIO(data))
    parts = []
    for slide in prs.slides:
        for shape in slide.shapes:
            if hasattr(shape, "text") and shape.text.strip():
                parts.append(shape.text)
    return "\n".join(parts)


def _parse_xlsx(data: bytes) -> List[Dict]:
    """
    Excel formati:
      A = savol | B,C,D,E = variantlar | F = to'g'ri javob (A/B/C/D yoki 1/2/3/4)
    """
    import openpyxl
    wb  = openpyxl.load_workbook(io.BytesIO(data), read_only=True, data_only=True)
    ws  = wb.active
    qs: List[Dict] = []

    for row in ws.iter_rows(min_row=2, values_only=True):
        if not row or not row[0]:
            continue
        question = str(row[0]).strip()
        options  = [str(row[i]).strip() for i in range(1, 5)
                    if i < len(row) and row[i]]
        if len(options) < 2:
            continue
        correct_id = 0
        if len(row) >= 6 and row[5]:
            mapping    = {"A": 0, "B": 1, "C": 2, "D": 3,
                          "1": 0, "2": 1, "3": 2, "4": 3}
            correct_id = mapping.get(str(row[5]).strip().upper(), 0)
        qs.append({
            "question":          question,
            "options":           options[:4],
            "correct_option_id": correct_id,
        })

    return [shuffle_options(q) for q in qs]


def _parse_pdf(data: bytes) -> str:
    import fitz
    doc = fitz.open(stream=data, filetype="pdf")
    return "\n".join(page.get_text() for page in doc)


def _parse_image(data: bytes) -> List[Dict]:
    """
    Rasmdan quiz savollarini o'qish — EasyOCR ishlatadi.
    Har xil formatdagi test rasmlarini tushunadi.
    """
    import easyocr
    import numpy as np
    from PIL import Image

    img     = Image.open(io.BytesIO(data)).convert("RGB")
    img_arr = np.array(img)

    reader  = easyocr.Reader(["ru", "en"], gpu=False, verbose=False)
    # detail=1 => har bir qatorni alohida o'qiymiz (bbox, text, confidence)
    results = reader.readtext(img_arr, detail=1, paragraph=False)

    # Faqat matnlarni olamiz, confidence >= 0.2
    lines = [r[1].strip() for r in results if r[2] >= 0.2 and r[1].strip()]

    return _parse_image_lines(lines)


def _parse_image_lines(lines: List[str]) -> List[Dict]:
    """
    EasyOCR dan kelgan qatorlar ro'yxatidan savollarni ajratadi.

    Quyidagi ko'rinishlarni tushunadi:
      - "# To'g'ri javob" yoki "# ..." — to'g'ri javob
      - "1." yoki "2." bilan boshlanuvchi — savol
      - "Variant 2", "Variant 3" — noto'g'ri javoblar
      - "Test Savollari", "===", "---" kabi sarlavha/chiziqlar — o'tkazib yuboriladi
    """
    questions: List[Dict] = []

    # Keraksiz qatorlarni filtrlash
    skip_patterns = [
        r"^test\s+savollari?$",
        r"^={2,}$",
        r"^-{2,}$",
        r"^\-+$",
        r"^[#\-=_\.]{3,}$",
        r"^\d+\s*[-–]\s*\d+$",   # sahifa raqamlari
    ]

    def should_skip(line: str) -> bool:
        l = line.strip().lower()
        if not l:
            return True
        for pat in skip_patterns:
            if re.match(pat, l, re.IGNORECASE):
                return True
        return False

    # Savol raqamini aniqlash
    q_re     = re.compile(r"^(\d+)[.)]\s*(.+)", re.UNICODE)
    # "# To'g'ri javob" yoki "# Toshkent" kabi
    hash_re  = re.compile(r"^#+\s*(.+)", re.UNICODE)

    current_q: str | None    = None
    current_opts: List[str]  = []
    correct_answer: str | None = None

    def save_question():
        nonlocal current_q, current_opts, correct_answer
        if current_q and correct_answer and len(current_opts) >= 1:
            # To'g'ri javob + noto'g'ri javoblar = options
            all_opts   = [correct_answer] + current_opts
            correct_id = 0  # to'g'ri javob 0-indeksda, shuffle_options aralashtiradi
            questions.append({
                "question":          current_q,
                "options":           all_opts[:4],
                "correct_option_id": correct_id,
            })
        current_q      = None
        current_opts   = []
        correct_answer = None

    for line in lines:
        line = line.strip()
        if should_skip(line):
            continue

        # Yangi savol boshlanganda oldingi saqlash
        qm = q_re.match(line)
        if qm:
            save_question()
            current_q      = qm.group(2).strip()
            current_q      = re.sub(r"^[#\s]+", "", current_q).strip()
            current_opts   = []
            correct_answer = None
            continue

        # # bilan boshlanuvchi = to'g'ri javob
        hm = hash_re.match(line)
        if hm:
            correct_answer = hm.group(1).strip()
            continue

        # Qolgan qatorlar = noto'g'ri variant (agar savol boshlangan bo'lsa)
        if current_q is not None and correct_answer is not None:
            # Raqamli "Variant 2" kabi yoki oddiy matn
            clean = re.sub(r"^[Vv]ariant\s*\d+\s*", "", line).strip()
            opt   = clean if clean else line
            if opt and len(current_opts) < 3:
                current_opts.append(opt)

    save_question()  # Oxirgi savolni saqlash

    return [shuffle_options(q) for q in questions]


# ──────────────────────────────────────────────────────────
#  Asosiy funksiya — tashqi chaqiruv
# ──────────────────────────────────────────────────────────
async def parse_file(
    file_bytes: bytes,
    file_name: str,
) -> Tuple[List[Dict], str]:
    """
    Faylni o'qib savollar ro'yxatini qaytaradi.
    Returns: (questions_list, error_message)
    error_message bo'sh bo'lsa — muvaffaqiyatli.
    """
    ext = file_name.rsplit(".", 1)[-1].lower() if "." in file_name else ""

    try:
        if ext == "txt":
            text = file_bytes.decode("utf-8", errors="ignore")
            qs   = parse_questions_from_text(text)

        elif ext in ("doc", "docx"):
            text = _parse_docx(file_bytes)
            qs   = parse_questions_from_text(text)

        elif ext == "pptx":
            text = _parse_pptx(file_bytes)
            qs   = parse_questions_from_text(text)

        elif ext in ("xls", "xlsx"):
            qs = _parse_xlsx(file_bytes)

        elif ext == "pdf":
            text = _parse_pdf(file_bytes)
            qs   = parse_questions_from_text(text)

        elif ext in ("jpg", "jpeg", "png", "bmp", "webp", "tiff"):
            # _parse_image endi to'g'ridan-to'g'ri List[Dict] qaytaradi
            qs = _parse_image(file_bytes)

        else:
            try:
                text = file_bytes.decode("utf-8", errors="ignore")
                qs   = parse_questions_from_text(text)
            except Exception:
                return [], (
                    "❌ Bu fayl formati qo'llab-quvvatlanmaydi.\n"
                    "📋 Formatlar: .txt  .docx  .pptx  .xlsx  .pdf  rasm"
                )

    except ImportError as exc:
        lib = str(exc).split("'")[1] if "'" in str(exc) else str(exc)
        log.error(f"ImportError: {exc}")
        return [], f"❌ Kutubxona o'rnatilmagan: *{lib}*\n(Xatolik: {exc})"

    except Exception as exc:
        log.exception("Fayl o'qishda xatolik")
        return [], f"❌ Fayl o'qishda xatolik: {exc}"

    if not qs:
        return [], (
            "❌ Fayldan savollar topilmadi!\n\n"
            "📝 *Format 1 (==== ajratuvchi):*\n"
            "```\n"
            "Savol matni?\n"
            "====\n"
            "# To'g'ri javob\n"
            "====\n"
            "Variant 2\n"
            "====\n"
            "Variant 3\n"
            "====\n"
            "Variant 4\n"
            "```\n\n"
            "📝 *Format 2 (klassik):*\n"
            "```\n"
            "1. Savol matni?\n"
            "A) Variant 1\n"
            "B) Variant 2\n"
            "C) To'g'ri variant\n"
            "D) Variant 4\n"
            "Javob: C\n"
            "```"
        )

    # Telegram cheklovi: savol ≤ 300 belgi, variant ≤ 100 belgi
    cleaned: List[Dict] = []
    for q in qs:
        opts = [o[:95] for o in q["options"]]
        if len(opts) < 2:
            continue
        cleaned.append({
            "question":          q["question"][:295],
            "options":           opts,
            "correct_option_id": min(q["correct_option_id"], len(opts) - 1),
        })

    if not cleaned:
        return [], "❌ Savollar Telegram chekloviga to'g'ri kelmadi."

    return cleaned, ""
